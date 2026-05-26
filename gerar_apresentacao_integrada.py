"""
Apresentação integrada III MCSM 2026
Combina o tema/estrutura da versão Gamma com as figuras reais de análise.
15 slides.
"""

import os, numpy as np, matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PASTA   = r"C:\Users\Jordana\OneDrive - FUNAI - Fundação Nacional dos Povos Indígenas\projetos programação\sertão"
FIGURAS = os.path.join(PASTA, "figuras")

# ── Paleta ─────────────────────────────────────────────────────
NAVY   = RGBColor(0x1B, 0x35, 0x66)
AZUL   = RGBColor(0x21, 0x76, 0xAE)
CREAM  = RGBColor(0xF5, 0xF0, 0xE8)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
CINZA  = RGBColor(0x44, 0x44, 0x44)
OURO   = RGBColor(0xC8, 0x9B, 0x3A)

# ══════════════════════════════════════════════════════════════
# FIGURAS AUXILIARES
# ══════════════════════════════════════════════════════════════

def gerar_diagrama_mlp():
    fig, ax = plt.subplots(figsize=(7, 5.5))
    fig.patch.set_facecolor("#F5F0E8")
    ax.set_facecolor("#F5F0E8")
    ax.set_xlim(0, 10); ax.set_ylim(0, 10)
    ax.axis("off")

    camadas = [
        (1.2, [8.5, 6.5, 4.5],          "#5B9BD5", "Entrada\n(3)"),
        (3.8, [9.0,7.5,6.0,4.5,3.0,1.5],"#E07A5F", "Oculta 1\n(16 — tanh)"),
        (6.4, [8.5,7.0,5.5,4.0,2.5,1.0],"#6C5B7B", "Oculta 2\n(8 — tanh)"),
        (9.0, [5.0],                     "#43AA8B", "Saída\n(linear)"),
    ]
    radii = [0.38, 0.30, 0.30, 0.42]

    # conexões
    for ci in range(len(camadas)-1):
        x0, ys0, _, _ = camadas[ci]
        x1, ys1, _, _ = camadas[ci+1]
        for y0 in ys0:
            for y1 in ys1:
                ax.plot([x0+radii[ci], x1-radii[ci+1]], [y0, y1],
                        color="#AAAAAA", lw=0.4, alpha=0.5, zorder=1)

    # neurônios
    for ci, (x, ys, cor, rot) in enumerate(camadas):
        r = radii[ci]
        for y in ys:
            c = plt.Circle((x, y), r, color=cor, zorder=3, ec="white", lw=1.2)
            ax.add_patch(c)
            ax.text(x, y, "σ", ha="center", va="center",
                    fontsize=9, color="white", fontweight="bold", zorder=4)
        ax.text(x, -0.3, rot, ha="center", va="top", fontsize=8.5,
                color="#1B3566", fontweight="bold")

    # setas entrada/saída
    ax.annotate("", xy=(0.82, 5.5), xytext=(0.1, 5.5),
                arrowprops=dict(arrowstyle="->", color="#1B3566", lw=1.5))
    ax.text(0.05, 5.5, "fc\nρ\nd", ha="right", va="center",
            fontsize=9, color="#1B3566", fontweight="bold")
    ax.annotate("", xy=(9.9, 5.0), xytext=(9.42, 5.0),
                arrowprops=dict(arrowstyle="->", color="#43AA8B", lw=1.5))
    ax.text(9.95, 5.0, "P\n(kN)", ha="left", va="center",
            fontsize=9, color="#1B3566", fontweight="bold")

    ax.set_title("Arquitetura MLP — 3 → 16 → 8 → 1",
                 fontsize=12, color="#1B3566", fontweight="bold", pad=10)
    plt.tight_layout()
    out = os.path.join(FIGURAS, "fig_mlp_arquitetura.png")
    plt.savefig(out, dpi=160, bbox_inches="tight", facecolor="#F5F0E8")
    plt.close()
    print("Diagrama MLP salvo.")
    return out


def gerar_grafico_metricas():
    fig, ax = plt.subplots(figsize=(9, 5))
    fig.patch.set_facecolor("#F5F0E8")
    ax.set_facecolor("#F5F0E8")

    metricas = ["MSE (kN²)", "RMSE (kN)", "MAE (kN)", "R² (Score)"]
    valores  = [0.292,       0.540,       0.437,      0.9998]
    cores    = ["#2176AE","#2176AE","#2176AE","#1B3566"]

    bars = ax.bar(metricas, valores, color=cores, width=0.5, zorder=3)
    for bar, val in zip(bars, valores):
        ax.text(bar.get_x() + bar.get_width()/2,
                bar.get_height() + 0.015,
                f"{val:.4f}" if val < 1 else f"{val:.3f}",
                ha="center", va="bottom", fontsize=13,
                color="#1B3566", fontweight="bold")

    ax.set_ylim(0, 1.18)
    ax.spines[["top","right","left"]].set_visible(False)
    ax.yaxis.set_visible(False)
    ax.tick_params(axis="x", labelsize=13, colors="#1B3566")
    ax.set_facecolor("#F5F0E8")
    for spine in ax.spines.values():
        spine.set_color("#CCCCCC")
    ax.xaxis.set_tick_params(length=0)
    ax.grid(axis="y", linestyle="--", alpha=0.3, zorder=0)
    ax.set_title("Métricas de Desempenho — MLP (conjunto de treinamento)",
                 fontsize=13, color="#1B3566", fontweight="bold", pad=12)
    plt.tight_layout()
    out = os.path.join(FIGURAS, "fig_metricas_barras.png")
    plt.savefig(out, dpi=160, bbox_inches="tight", facecolor="#F5F0E8")
    plt.close()
    print("Gráfico de métricas salvo.")
    return out


fig_mlp = gerar_diagrama_mlp()
fig_met = gerar_grafico_metricas()

# ══════════════════════════════════════════════════════════════
# HELPERS PPTX
# ══════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def bg_cor(slide, cor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = cor

def rect(slide, l, t, w, h, cor, linha=False):
    s = slide.shapes.add_shape(1,
        Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = cor
    if linha:
        s.line.color.rgb = cor
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def txt(slide, texto, l, t, w, h,
        sz=18, bold=False, italic=False,
        cor=NAVY, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = texto
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = cor
    return tb

def bullets(slide, itens, l, t, w, h,
            sz=17, cor_txt=NAVY, cor_dot=AZUL):
    tb = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(itens):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5)
        r0 = p.add_run(); r0.text = "●  "
        r0.font.size = Pt(sz); r0.font.color.rgb = cor_dot
        r1 = p.add_run(); r1.text = item
        r1.font.size = Pt(sz); r1.font.color.rgb = cor_txt

def img(slide, path, l, t, w=None, h=None):
    if not os.path.exists(path): return
    kw = {}
    if w: kw["width"]  = Inches(w)
    if h: kw["height"] = Inches(h)
    slide.shapes.add_picture(path, Inches(l), Inches(t), **kw)

def cabecalho(slide, titulo, subtitulo=None):
    """Faixa navy no topo + barra lateral esquerda."""
    rect(slide, 0, 0, 13.33, 1.25, NAVY)
    rect(slide, 0, 0, 0.12,  7.50, AZUL)
    txt(slide, titulo, 0.25, 0.1, 12.8, 1.0,
        sz=30, bold=True, cor=BRANCO)
    if subtitulo:
        txt(slide, subtitulo, 0.25, 0.78, 12.8, 0.45,
            sz=14, italic=True, cor=RGBColor(0xC0,0xD8,0xF0))

def rodape(slide, texto="III MCSM 2026  |  UNIMONTES"):
    rect(slide, 0, 7.05, 13.33, 0.45, NAVY)
    txt(slide, texto, 0.3, 7.08, 12.7, 0.38,
        sz=11, italic=True, cor=BRANCO, align=PP_ALIGN.CENTER)

def card(slide, l, t, w, h, titulo, corpo, sz_tit=15, sz_corp=13):
    rect(slide, l, t, w, h, BRANCO, linha=True)
    rect(slide, l, t+h-0.08, w, 0.08, NAVY)   # barra inferior
    txt(slide, titulo, l+0.15, t+0.18, w-0.3, 0.45,
        sz=sz_tit, bold=True, cor=AZUL)
    txt(slide, corpo, l+0.15, t+0.65, w-0.3, h-0.8,
        sz=sz_corp, cor=CINZA)

# ══════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_cor(sl, NAVY)
rect(sl, 0, 0, 0.18, 7.5, AZUL)
rect(sl, 0, 6.55, 13.33, 0.95, AZUL)

txt(sl, "Previsão da Carga de Ruptura em", 0.4, 0.8, 12.6, 1.1,
    sz=40, bold=True, cor=BRANCO, align=PP_ALIGN.CENTER)
txt(sl, "Vigas de Concreto Armado", 0.4, 1.85, 12.6, 1.0,
    sz=40, bold=True, cor=BRANCO, align=PP_ALIGN.CENTER)

rect(sl, 3.5, 3.05, 6.3, 0.06, AZUL)

txt(sl, "Utilizando Redes Neurais Multilayer Perceptron (MLP)",
    0.4, 3.2, 12.6, 0.6,
    sz=18, italic=True, cor=CREAM, align=PP_ALIGN.CENTER)

txt(sl, "Maria Luiza A. Xavier   |   Jordana G. Fernandes   |   Rayre J. V. Silva",
    0.4, 4.2, 12.6, 0.5,
    sz=15, cor=RGBColor(0xC8,0xB8,0x96), align=PP_ALIGN.CENTER)
txt(sl, "Universidade Estadual de Montes Claros (UNIMONTES)",
    0.4, 4.75, 12.6, 0.45,
    sz=13, italic=True, cor=RGBColor(0xA0,0xB4,0xC8), align=PP_ALIGN.CENTER)

txt(sl, "III MCSM — 13 a 15 de agosto de 2026",
    0.4, 6.6, 8.0, 0.38,
    sz=12, italic=True, cor=BRANCO)
txt(sl, "github.com/JGF1987/mlp-viga-concreto-mcsm",
    6.5, 6.6, 6.5, 0.38,
    sz=11, cor=RGBColor(0xA0,0xC8,0xE8), align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════
# SLIDE 2 — DESAFIOS NO DIMENSIONAMENTO
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_cor(sl, CREAM)
cabecalho(sl, "DESAFIOS NO DIMENSIONAMENTO")

txt(sl, "Limitações Normativas", 0.4, 1.45, 7.0, 0.55,
    sz=20, bold=True, cor=AZUL)
txt(sl, "Equações das normas (NBR 6118, ACI 318) são calibradas para\n"
        "configurações típicas e utilizam hipóteses simplificadoras.",
    0.4, 2.0, 6.8, 0.8, sz=15, cor=CINZA)

bullets(sl, [
    "Negligenciam interações não lineares complexas",
    "Limitadas em geometrias e taxas de armadura não convencionais",
    "Dificuldade em capturar comportamentos mecânicos intrincados",
    "RMSE empírico típico: 12–17 kN·m (Khan et al., 2023)",
], 0.4, 2.85, 6.8, 3.2, sz=16)

# Caixa direita — contexto e motivação
rect(sl, 7.4, 1.4, 5.6, 5.2, BRANCO, linha=True)
rect(sl, 7.4, 6.52, 5.6, 0.08, NAVY)
txt(sl, "Por que Machine Learning?", 7.6, 1.55, 5.2, 0.5,
    sz=15, bold=True, cor=AZUL)
bullets(sl, [
    "Modelos de ML superam equações empíricas",
    "Capturam relações não lineares automaticamente",
    "Aplicáveis a ensaios não destrutivos",
    "Escaláveis para bases de dados reais",
], 7.55, 2.1, 5.2, 4.3, sz=14, cor_txt=CINZA, cor_dot=NAVY)

rodape(sl)

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — APRENDIZADO DE MÁQUINA (3 cards)
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_cor(sl, CREAM)
cabecalho(sl, "APRENDIZADO DE MÁQUINA")

for i, (tit, corp) in enumerate([
    ("Precisão Elevada",
     "Modelos de ML superam consistentemente expressões normativas em tarefas de previsão estrutural."),
    ("Relacionamentos Ocultos",
     "Capacidade de mapear relações não lineares entre resistência, taxa de armadura e geometria."),
    ("Eficiência de Tempo",
     "Métodos promissores para análise rápida e ensaios não destrutivos em estruturas reais."),
]):
    card(sl, 0.35 + i*4.35, 1.5, 4.1, 5.2, tit, corp, sz_tit=17, sz_corp=14)

rodape(sl)

# ══════════════════════════════════════════════════════════════
# SLIDE 4 — FUNDAMENTOS DA REDE MLP
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_cor(sl, CREAM)
cabecalho(sl, "FUNDAMENTOS DA REDE MLP")

img(sl, fig_mlp, l=0.3, t=1.35, w=6.8)

txt(sl, "Arquitetura Multicamadas", 7.4, 1.5, 5.6, 0.55,
    sz=20, bold=True, cor=AZUL)
txt(sl, "Processamento através de neurônios artificiais\ntotalmente conectados:",
    7.4, 2.1, 5.6, 0.7, sz=14, cor=CINZA)
bullets(sl, [
    "Camada de Entrada: recebe os parâmetros físicos (fc, ρ, d)",
    "Camadas Ocultas: aprendem representações intermediárias",
    "Camada de Saída: estimativa contínua de P (kN)",
    "Pesos ajustados por retropropagação com Adam",
], 7.4, 2.9, 5.6, 3.5, sz=14)

rodape(sl)

# ══════════════════════════════════════════════════════════════
# SLIDE 5 — ATIVAÇÃO E OTIMIZAÇÃO (3 cards)
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_cor(sl, CREAM)
cabecalho(sl, "ATIVAÇÃO E OTIMIZAÇÃO")

for i, (tit, corp) in enumerate([
    ("Tangente Hiperbólica (tanh)",
     "Mapeia saídas para (−1, 1), produzindo gradientes suaves e evitando saturação assimétrica. Preferível à ReLU em regressões com escalas heterogêneas."),
    ("Otimizador Adam",
     "Combina momentum de 1ª e 2ª ordem, adaptando a taxa de aprendizado individualmente para cada parâmetro — convergência mais rápida."),
    ("Normalização Min-Max",
     "Escala variáveis para [0, 1]: x' = (x − xₘᵢₙ)/(xₘₐₓ − xₘᵢₙ). Evita saturação precoce dos neurônios durante o treinamento."),
]):
    card(sl, 0.35 + i*4.35, 1.5, 4.1, 5.2, tit, corp, sz_tit=15, sz_corp=13)

rodape(sl)

# ══════════════════════════════════════════════════════════════
# SLIDE 6 — GERAÇÃO DE DADOS SINTÉTICOS
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_cor(sl, CREAM)
cabecalho(sl, "GERAÇÃO DE DADOS SINTÉTICOS")

txt(sl,
    "A base de dados foi gerada por uma relação analítica que introduz\n"
    "não linearidade mecânica real e ruído gaussiano experimental:",
    0.4, 1.45, 12.6, 0.75, sz=15, cor=CINZA)

rect(sl, 1.0, 2.35, 11.3, 1.4, BRANCO, linha=True)
txt(sl, "P  =  0,35 · fc  +  12 · ρ  +  2,1 · d  −  0,02 · fc · d  +  ε",
    1.2, 2.55, 10.9, 0.95,
    sz=26, bold=True, cor=NAVY, align=PP_ALIGN.CENTER)

bullets(sl, [
    "fc ∈ [25, 50] MPa — Resistência à compressão do concreto",
    "ρ  ∈ [0,5%, 2,5%] — Taxa de armadura",
    "d  ∈ [30, 60] cm — Altura útil da seção transversal",
    "ε  ~ N(0; 1,5) — Ruído gaussiano (simula variações experimentais)",
    "Termo −0,02 · fc · d: interação não linear intr ínseca ao comportamento estrutural",
], 0.5, 3.95, 12.4, 2.8, sz=15)

rodape(sl)

# ══════════════════════════════════════════════════════════════
# SLIDE 7 — AMOSTRAS DA BASE DE DADOS (tabela completa)
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_cor(sl, CREAM)
cabecalho(sl, "AMOSTRAS DA BASE DE DADOS")

dados = [
    ("fc (MPa)", "ρ (%)", "d (cm)", "P Real (kN)", "P Previsto (kN)", "Erro (%)"),
    ("28","0,8","35","78,4","78,26","−0,17"),
    ("32","1,0","38","91,2","90,99","−0,23"),
    ("35","1,2","42","104,5","105,04","+0,51"),
    ("38","1,5","45","119,3","120,33","+0,86"),
    ("42","1,8","48","138,7","138,84","+0,10"),
    ("45","2,0","52","156,1","156,15","+0,03"),
    ("48","2,2","55","172,9","172,14","−0,44"),
    ("50","2,5","60","191,4","191,81","+0,22"),
    ("30","2,0","32","93,8","93,52","−0,30"),
    ("40","1,3","50","128,6","127,79","−0,63"),
]
col_w = [1.7, 1.25, 1.25, 1.85, 1.95, 1.45]
tbl = sl.shapes.add_table(
    len(dados), 6,
    Inches(0.4), Inches(1.45),
    Inches(12.55), Inches(5.7)).table
for ci, w in enumerate(col_w):
    tbl.columns[ci].width = Inches(w)
for ri, row in enumerate(dados):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.size = Pt(13 if ri > 0 else 12)
        run.font.bold = (ri == 0)
        run.font.color.rgb = BRANCO if ri == 0 else CINZA
        cell.fill.solid()
        if ri == 0:
            cell.fill.fore_color.rgb = NAVY
        elif ri % 2 == 0:
            cell.fill.fore_color.rgb = RGBColor(0xE8,0xF0,0xF8)
        else:
            cell.fill.fore_color.rgb = BRANCO

rodape(sl, "Dados normalizados min-max para o treinamento  |  III MCSM 2026")

# ══════════════════════════════════════════════════════════════
# SLIDE 8 — CONFIGURAÇÃO DA REDE MLP
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_cor(sl, CREAM)
cabecalho(sl, "CONFIGURAÇÃO DA REDE MLP")

bullets(sl, [
    "Entrada: 3 neurônios (fc, ρ, d)",
    "Camada Oculta 1: 16 neurônios (ativação Tanh)",
    "Camada Oculta 2: 8 neurônios (ativação Tanh)",
    "Saída: 1 neurônio (ativação Linear — Regressão)",
    "Otimizador: Adam  |  η = 0,01",
    "Épocas: 1.000  |  Função de custo: MSE",
], 0.4, 1.4, 6.5, 5.2, sz=17)

img(sl, fig_mlp, l=7.1, t=1.4, w=6.0)

rodape(sl)

# ══════════════════════════════════════════════════════════════
# SLIDE 9 — RESULTADOS DE TREINAMENTO (gráfico de barras)
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_cor(sl, CREAM)
cabecalho(sl, "RESULTADOS DE TREINAMENTO")

img(sl, fig_met, l=1.5, t=1.35, w=10.3)

rodape(sl)

# ══════════════════════════════════════════════════════════════
# SLIDE 10 — CURVA DE TREINAMENTO + DISPERSÃO
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_cor(sl, CREAM)
cabecalho(sl, "CURVA DE TREINAMENTO E DISPERSÃO")

img(sl, os.path.join(FIGURAS,"fig1_curva_treinamento.png"),  l=0.3,  t=1.35, w=6.35)
img(sl, os.path.join(FIGURAS,"fig2_real_vs_previsto.png"),   l=6.85, t=1.35, w=6.35)

txt(sl, "MSE cai ~4 ordens de grandeza em 1.000 épocas — convergência estável",
    0.3, 6.55, 6.35, 0.42, sz=12, italic=True, cor=CINZA, align=PP_ALIGN.CENTER)
txt(sl, "MLP (R²=0,9998) supera amplamente Regressão Linear (R²=0,9960)",
    6.85, 6.55, 6.35, 0.42, sz=12, italic=True, cor=CINZA, align=PP_ALIGN.CENTER)

rodape(sl)

# ══════════════════════════════════════════════════════════════
# SLIDE 11 — ANÁLISE DE SENSIBILIDADE + SUPERFÍCIE 3D
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_cor(sl, CREAM)
cabecalho(sl, "SENSIBILIDADE E SUPERFÍCIE 3D")

img(sl, os.path.join(FIGURAS,"fig3_sensibilidade.png"), l=0.25, t=1.35, w=12.85)
img(sl, os.path.join(FIGURAS,"fig4_superficie_3d.png"), l=0.25, t=4.3,  w=12.85)

rodape(sl)

# ══════════════════════════════════════════════════════════════
# SLIDE 12 — ANÁLISE DE ADERÊNCIA
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_cor(sl, CREAM)
cabecalho(sl, "ANÁLISE DE ADERÊNCIA")

txt(sl, "Captura da Não Linearidade", 0.4, 1.5, 7.0, 0.55,
    sz=20, bold=True, cor=AZUL)
txt(sl,
    "A aderência à linha de identidade confirma que a MLP\n"
    "aprendeu com precisão o termo de interação fc·d.",
    0.4, 2.15, 6.8, 0.85, sz=15, cor=CINZA)
txt(sl,
    "Erros percentuais máximos inferiores a 1% em todas as\n"
    "10 amostras validam a eficácia da arquitetura proposta\n"
    "frente a modelos tradicionais.",
    0.4, 3.1, 6.8, 1.0, sz=15, cor=CINZA)

bullets(sl, [
    "R² = 0,9998  vs.  R² = 0,9960 (Regressão Linear)",
    "RMSE = 0,54 kN  vs.  RMSE ≈ 12–17 kN·m (modelos empíricos)",
    "Superior ao melhor modelo da literatura: R² = 0,982\n"
    "  (Özyüksel Çiftçioğlu et al., 2025)",
], 0.4, 4.2, 6.8, 2.5, sz=14)

img(sl, os.path.join(FIGURAS,"fig5_residuos.png"), l=7.2, t=1.4, w=5.9)

rodape(sl)

# ══════════════════════════════════════════════════════════════
# SLIDE 13 — CONCLUSÕES (3 cards)
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_cor(sl, CREAM)
cabecalho(sl, "CONCLUSÕES DO ESTUDO")

for i, (tit, corp) in enumerate([
    ("Validação",
     "A MLP demonstrou alta capacidade de ajuste (R²=0,9998), capturando efeitos não lineares intrínsecos ao concreto armado."),
    ("Vantagem Técnica",
     "Redução drástica do erro comparado a modelos empíricos (RMSE de ~15 kN reduzido para 0,54 kN) e superiora à literatura recente."),
    ("Perspectivas",
     "Ampliar a base com dados experimentais reais e avaliar a generalização via validação cruzada k-fold."),
]):
    card(sl, 0.35 + i*4.35, 1.5, 4.1, 5.2, tit, corp, sz_tit=17, sz_corp=13)

rodape(sl)

# ══════════════════════════════════════════════════════════════
# SLIDE 14 — DISPONIBILIDADE + AGRADECIMENTOS
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_cor(sl, NAVY)
rect(sl, 0, 0, 0.18, 7.5, AZUL)

txt(sl, "Disponibilidade do Código", 0.5, 0.7, 12.5, 0.65,
    sz=26, bold=True, cor=BRANCO)
rect(sl, 0.4, 1.45, 12.5, 1.0, AZUL)
txt(sl, "github.com/JGF1987/mlp-viga-concreto-mcsm",
    0.5, 1.55, 12.4, 0.8,
    sz=22, bold=True, cor=BRANCO, align=PP_ALIGN.CENTER)
txt(sl,
    "problema4_mlp.py   |   visualizacoes_p4.py   |   gerar_word_artigo.py\n"
    "figuras/ (6 figuras em alta resolução)   |   artigo_problema4.tex   |   artigo_problema4.docx",
    0.5, 2.6, 12.5, 0.9, sz=13, cor=CREAM, align=PP_ALIGN.CENTER)

rect(sl, 0.4, 3.7, 12.5, 0.06, AZUL)

txt(sl, "Agradecimentos", 0.5, 3.9, 12.5, 0.55,
    sz=22, bold=True, cor=BRANCO)
txt(sl, "Professor Marcos e Universidade Estadual de Montes Claros — UNIMONTES",
    0.5, 4.5, 12.5, 0.55, sz=18, italic=True, cor=CREAM)

rect(sl, 0.4, 5.25, 12.5, 0.06, AZUL)

txt(sl,
    "Maria Luiza Almeida Xavier  ·  maria.luiza.a.xavier@gmail.com\n"
    "Jordana Gabriela Fernandes   ·  jordana.funai@gmail.com\n"
    "Rayre Janaína Vieira Silva    ·  rayrejanaina2@gmail.com",
    0.5, 5.4, 12.5, 1.4, sz=14, cor=CREAM)

# ══════════════════════════════════════════════════════════════
# SLIDE 15 — PERGUNTAS & RESPOSTAS
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_cor(sl, CREAM)
rect(sl, 0, 0, 13.33, 0.12, NAVY)
rect(sl, 0, 7.38, 13.33, 0.12, NAVY)

txt(sl, "Perguntas & Respostas",
    0.4, 2.1, 12.6, 1.3,
    sz=46, bold=True, cor=NAVY, align=PP_ALIGN.CENTER)
rect(sl, 5.5, 3.55, 2.3, 0.08, AZUL)
txt(sl, "Agradecemos pela atenção!",
    0.4, 3.8, 12.6, 0.55,
    sz=18, italic=True, cor=CINZA, align=PP_ALIGN.CENTER)

txt(sl, "Códigos disponíveis em:",
    0.4, 4.9, 12.6, 0.45,
    sz=15, bold=True, cor=NAVY, align=PP_ALIGN.CENTER)
txt(sl, "github.com/JGF1987/mlp-viga-concreto-mcsm",
    0.4, 5.4, 12.6, 0.5,
    sz=16, cor=AZUL, align=PP_ALIGN.CENTER)

# ── Salvar ─────────────────────────────────────────────────────
out = os.path.join(PASTA, "apresentacao_mcsm2026_integrada.pptx")
prs.save(out)
print(f"\nApresentação integrada salva: {out}")
print(f"Total: {len(prs.slides)} slides")
print("Slides: Capa | Desafios | ML | Fund.MLP | Ativação | Dados | Tabela | "
      "Config | Métricas | Curva+Dispersão | Sensib+3D | Aderência | Conclusões | GitHub | Q&A")
