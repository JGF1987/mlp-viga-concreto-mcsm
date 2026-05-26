"""
Gera apresentação PowerPoint do artigo III MCSM 2026
Previsão de Carga de Ruptura em Vigas de Concreto Armado com MLP
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu
import os

PASTA = r"C:\Users\Jordana\OneDrive - FUNAI - Fundação Nacional dos Povos Indígenas\projetos programação\sertão"
FIGURAS = os.path.join(PASTA, "figuras")

# ── Paleta de cores ────────────────────────────────────────────
AZUL_ESC  = RGBColor(0x17, 0x39, 0x6B)   # azul escuro UNIMONTES
AZUL_MED  = RGBColor(0x21, 0x76, 0xAE)   # azul médio
AZUL_CLA  = RGBColor(0xD6, 0xE4, 0xF0)   # azul claro (fundo)
BRANCO    = RGBColor(0xFF, 0xFF, 0xFF)
CINZA     = RGBColor(0x55, 0x55, 0x55)
LARANJA   = RGBColor(0xE0, 0x7A, 0x5F)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # layout em branco

# ══════════════════════════════════════════════════════════════
# HELPERS
# ══════════════════════════════════════════════════════════════

def bg(slide, cor=AZUL_ESC):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = cor

def caixa(slide, texto, left, top, width, height,
          size=20, bold=False, cor_txt=BRANCO, cor_bg=None,
          italic=False, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = texto
    run.font.size    = Pt(size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = cor_txt
    if cor_bg:
        txb.fill.solid()
        txb.fill.fore_color.rgb = cor_bg
    return txb

def linha_h(slide, top, cor=AZUL_MED, espessura=Pt(2)):
    from pptx.util import Pt as Ptx
    ln = slide.shapes.add_connector(
        1,  # msoConnectorStraight
        Inches(0.4), Inches(top),
        Inches(12.93), Inches(top))
    ln.line.color.rgb = cor
    ln.line.width     = espessura

def figura(slide, caminho, left, top, width=None, height=None):
    if not os.path.exists(caminho):
        return
    if width and height:
        slide.shapes.add_picture(caminho,
            Inches(left), Inches(top), Inches(width), Inches(height))
    elif width:
        slide.shapes.add_picture(caminho,
            Inches(left), Inches(top), width=Inches(width))
    else:
        slide.shapes.add_picture(caminho,
            Inches(left), Inches(top), height=Inches(height))

def retangulo(slide, left, top, width, height, cor=AZUL_MED, radius=False):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shp = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = cor
    shp.line.fill.background()
    return shp

def bullet(slide, itens, left, top, width, height,
           size=18, cor_txt=BRANCO, cor_bullet=LARANJA, espaco=1.2):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(itens):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        # marcador
        r0 = p.add_run()
        r0.text = "▸  "
        r0.font.size  = Pt(size)
        r0.font.color.rgb = cor_bullet
        r0.font.bold  = True
        # texto
        r1 = p.add_run()
        r1.text = item
        r1.font.size  = Pt(size)
        r1.font.color.rgb = cor_txt

# ══════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, AZUL_ESC)

# Faixa lateral esquerda
retangulo(sl, 0, 0, 0.18, 7.5, AZUL_MED)

# Faixa inferior
retangulo(sl, 0, 6.8, 13.33, 0.7, AZUL_MED)

# Título principal
caixa(sl,
    "PREVISÃO DE CARGA DE RUPTURA\nEM VIGAS DE CONCRETO ARMADO\nCOM REDES NEURAIS MLP",
    0.6, 1.2, 12.4, 2.8,
    size=34, bold=True, cor_txt=BRANCO, align=PP_ALIGN.CENTER)

# Subtítulo: evento
caixa(sl,
    "III MCSM — Modelagem Computacional do Sertão Mineiro\n13–15 de agosto de 2026",
    0.6, 4.1, 12.4, 0.8,
    size=18, italic=True, cor_txt=AZUL_CLA, align=PP_ALIGN.CENTER)

# Autoras
caixa(sl,
    "Maria Luiza Almeida Xavier   |   Jordana Gabriela Fernandes   |   Rayre Janaína Vieira Silva",
    0.6, 5.1, 12.4, 0.5,
    size=15, cor_txt=BRANCO, align=PP_ALIGN.CENTER)

# Instituição
caixa(sl,
    "Universidade Estadual de Montes Claros — UNIMONTES",
    0.6, 5.7, 12.4, 0.4,
    size=13, italic=True, cor_txt=AZUL_CLA, align=PP_ALIGN.CENTER)

# Rodapé
caixa(sl, "Agradecimentos: professor Marcos e UNIMONTES",
    0.3, 6.85, 8, 0.4,
    size=11, italic=True, cor_txt=BRANCO)
caixa(sl, "github.com/JGF1987/mlp-viga-concreto-mcsm",
    9.0, 6.85, 4.0, 0.4,
    size=11, cor_txt=AZUL_CLA, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════
# SLIDE 2 — MOTIVAÇÃO E OBJETIVO
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, BRANCO)
retangulo(sl, 0, 0, 13.33, 1.15, AZUL_ESC)
caixa(sl, "Motivação e Objetivo", 0.4, 0.15, 12.5, 0.85,
      size=28, bold=True, cor_txt=BRANCO)
linha_h(sl, 1.25)

bullet(sl, [
    "Vigas de concreto armado apresentam comportamento não linear complexo",
    "Equações normativas (NBR 6118:2023) introduzem hipóteses simplificadoras",
    "Métodos empíricos têm RMSE de 12–17 kN·m — precisão limitada",
    "Modelos de ML superam consistentemente equações normativas (Khan et al., 2023)",
], 0.5, 1.4, 7.5, 3.5, size=19, cor_txt=CINZA, cor_bullet=AZUL_MED)

# Caixa objetivo
retangulo(sl, 0.4, 5.0, 12.5, 1.9, AZUL_CLA)
caixa(sl, "Objetivo",
      0.55, 5.05, 3.0, 0.4, size=16, bold=True, cor_txt=AZUL_ESC)
caixa(sl,
    "Aplicar uma rede MLP para estimar a carga de ruptura P (kN) de vigas\n"
    "de concreto armado a partir de fc, ρ e d — e comparar com regressão linear.",
    0.55, 5.45, 12.1, 1.3, size=18, cor_txt=AZUL_ESC)

retangulo(sl, 0, 6.9, 13.33, 0.6, AZUL_MED)
caixa(sl, "III MCSM 2026  |  UNIMONTES", 0.4, 6.93, 12.5, 0.4,
      size=11, italic=True, cor_txt=BRANCO, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — FUNDAMENTOS: MLP
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, BRANCO)
retangulo(sl, 0, 0, 13.33, 1.15, AZUL_ESC)
caixa(sl, "Rede Neural MLP — Fundamentos", 0.4, 0.15, 12.5, 0.85,
      size=28, bold=True, cor_txt=BRANCO)
linha_h(sl, 1.25)

bullet(sl, [
    "Neurônio: z = Σ wᵢⱼ · aᵢ + b  →  a = φ(z)",
    "Ativação tanh: mapeia saída a (−1, 1), gradientes suaves, evita saturação assimétrica",
    "Normalização min-max: x' = (x − xₘᵢₙ) / (xₘₐₓ − xₘᵢₙ)",
    "Otimizador Adam: adapta taxa de aprendizado por parâmetro — convergência rápida",
    "Função de custo: MSE = (1/n) Σ (yᵢ − ŷᵢ)²",
], 0.5, 1.4, 8.2, 4.5, size=18, cor_txt=CINZA, cor_bullet=AZUL_MED)

# Arquitetura resumida
retangulo(sl, 8.6, 1.4, 4.4, 4.5, AZUL_CLA)
caixa(sl, "Arquitetura Adotada", 8.75, 1.5, 4.1, 0.45,
      size=15, bold=True, cor_txt=AZUL_ESC)
bullet(sl, [
    "Entrada: 3 neurônios (fc, ρ, d)",
    "Camada 1: 16 neurônios — tanh",
    "Camada 2:   8 neurônios — tanh",
    "Saída: 1 neurônio — linear",
    "Adam  |  η = 0,01  |  1.000 épocas",
], 8.7, 2.0, 4.1, 3.7, size=15, cor_txt=AZUL_ESC, cor_bullet=LARANJA)

retangulo(sl, 0, 6.9, 13.33, 0.6, AZUL_MED)
caixa(sl, "III MCSM 2026  |  UNIMONTES", 0.4, 6.93, 12.5, 0.4,
      size=11, italic=True, cor_txt=BRANCO, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 4 — METODOLOGIA: DADOS
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, BRANCO)
retangulo(sl, 0, 0, 13.33, 1.15, AZUL_ESC)
caixa(sl, "Metodologia — Base de Dados Sintética", 0.4, 0.15, 12.5, 0.85,
      size=28, bold=True, cor_txt=BRANCO)
linha_h(sl, 1.25)

caixa(sl,
    "Equação geradora (com ruído gaussiano ε ~ N(0; 1,5) e interação não linear fc · d):",
    0.5, 1.35, 12.5, 0.5, size=17, cor_txt=CINZA)
caixa(sl,
    "P  =  0,35 · fc  +  12 · ρ  +  2,1 · d  −  0,02 · fc · d  +  ε",
    0.5, 1.85, 12.5, 0.65, size=22, bold=True, cor_txt=AZUL_ESC, align=PP_ALIGN.CENTER)

# Tabela de dados
dados = [
    ("fc (MPa)", "ρ (%)", "d (cm)", "P real (kN)", "P prev. (kN)", "Erro (%)"),
    ("28", "0,8", "35", "78,4",  "78,26",  "−0,17"),
    ("32", "1,0", "38", "91,2",  "90,99",  "−0,23"),
    ("35", "1,2", "42", "104,5", "105,04", "+0,51"),
    ("38", "1,5", "45", "119,3", "120,33", "+0,86"),
    ("42", "1,8", "48", "138,7", "138,84", "+0,10"),
    ("45", "2,0", "52", "156,1", "156,15", "+0,03"),
    ("48", "2,2", "55", "172,9", "172,14", "−0,44"),
    ("50", "2,5", "60", "191,4", "191,81", "+0,22"),
    ("30", "2,0", "32", "93,8",  "93,52",  "−0,30"),
    ("40", "1,3", "50", "128,6", "127,79", "−0,63"),
]

rows, cols = len(dados), len(dados[0])
tbl = sl.shapes.add_table(rows, cols,
    Inches(0.4), Inches(2.65), Inches(12.5), Inches(4.1)).table

col_w = [1.7, 1.3, 1.3, 1.9, 1.9, 1.4]
for ci, w in enumerate(col_w):
    tbl.columns[ci].width = Inches(w)

for ri, row in enumerate(dados):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.size = Pt(13)
        run.font.bold = (ri == 0)
        run.font.color.rgb = BRANCO if ri == 0 else CINZA
        cell.fill.solid()
        if ri == 0:
            cell.fill.fore_color.rgb = AZUL_ESC
        elif ri % 2 == 0:
            cell.fill.fore_color.rgb = AZUL_CLA
        else:
            cell.fill.fore_color.rgb = BRANCO

retangulo(sl, 0, 6.9, 13.33, 0.6, AZUL_MED)
caixa(sl, "III MCSM 2026  |  UNIMONTES", 0.4, 6.93, 12.5, 0.4,
      size=11, italic=True, cor_txt=BRANCO, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 5 — RESULTADOS: MÉTRICAS
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, BRANCO)
retangulo(sl, 0, 0, 13.33, 1.15, AZUL_ESC)
caixa(sl, "Resultados — Métricas de Desempenho", 0.4, 0.15, 12.5, 0.85,
      size=28, bold=True, cor_txt=BRANCO)
linha_h(sl, 1.25)

metricas = [
    ("R²",   "0,9998", AZUL_ESC),
    ("RMSE", "0,54 kN", AZUL_MED),
    ("MAE",  "0,44 kN", AZUL_MED),
    ("MSE",  "0,292 kN²", AZUL_MED),
]
for i, (nome, val, cor) in enumerate(metricas):
    x = 0.4 + i * 3.15
    retangulo(sl, x, 1.45, 2.9, 1.8, cor)
    caixa(sl, nome, x, 1.5, 2.9, 0.6,
          size=20, bold=True, cor_txt=BRANCO, align=PP_ALIGN.CENTER)
    caixa(sl, val, x, 2.1, 2.9, 0.9,
          size=28, bold=True, cor_txt=BRANCO, align=PP_ALIGN.CENTER)

bullet(sl, [
    "Todos os erros percentuais individuais abaixo de 1%",
    "MLP capturou o termo não linear fc · d com precisão",
    "Regressão Linear: R² = 0,9960 | RMSE = 2,24 kN — desempenho inferior",
    "Özyüksel Çiftçioğlu et al. (2025): R² = 0,982 — MLP proposta é superior",
    "Khan et al. (2023): RMSE empírico de 12–17 kN·m vs. 0,54 kN aqui",
], 0.5, 3.45, 12.5, 3.1, size=17, cor_txt=CINZA, cor_bullet=AZUL_MED)

retangulo(sl, 0, 6.9, 13.33, 0.6, AZUL_MED)
caixa(sl, "III MCSM 2026  |  UNIMONTES", 0.4, 6.93, 12.5, 0.4,
      size=11, italic=True, cor_txt=BRANCO, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 6 — RESULTADOS: CURVA DE TREINAMENTO + DISPERSÃO
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, BRANCO)
retangulo(sl, 0, 0, 13.33, 1.15, AZUL_ESC)
caixa(sl, "Resultados — Curva de Treinamento e Dispersão", 0.4, 0.15, 12.5, 0.85,
      size=28, bold=True, cor_txt=BRANCO)
linha_h(sl, 1.25)

figura(sl, os.path.join(FIGURAS, "fig1_curva_treinamento.png"),
       left=0.3, top=1.35, width=6.3)
figura(sl, os.path.join(FIGURAS, "fig2_real_vs_previsto.png"),
       left=6.8, top=1.35, width=6.3)

caixa(sl, "MSE cai ~4 ordens de grandeza em 1.000 épocas",
      0.3, 6.35, 6.3, 0.45, size=13, italic=True, cor_txt=CINZA, align=PP_ALIGN.CENTER)
caixa(sl, "MLP (R²=0,9998) muito superior à Regressão Linear (R²=0,9960)",
      6.8, 6.35, 6.3, 0.45, size=13, italic=True, cor_txt=CINZA, align=PP_ALIGN.CENTER)

retangulo(sl, 0, 6.9, 13.33, 0.6, AZUL_MED)
caixa(sl, "III MCSM 2026  |  UNIMONTES", 0.4, 6.93, 12.5, 0.4,
      size=11, italic=True, cor_txt=BRANCO, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 7 — RESULTADOS: SENSIBILIDADE E SUPERFÍCIE 3D
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, BRANCO)
retangulo(sl, 0, 0, 13.33, 1.15, AZUL_ESC)
caixa(sl, "Resultados — Sensibilidade e Superfície 3D", 0.4, 0.15, 12.5, 0.85,
      size=28, bold=True, cor_txt=BRANCO)
linha_h(sl, 1.25)

figura(sl, os.path.join(FIGURAS, "fig3_sensibilidade.png"),
       left=0.3, top=1.35, width=12.7)
figura(sl, os.path.join(FIGURAS, "fig4_superficie_3d.png"),
       left=0.3, top=4.55, width=12.7)

retangulo(sl, 0, 6.9, 13.33, 0.6, AZUL_MED)
caixa(sl, "III MCSM 2026  |  UNIMONTES", 0.4, 6.93, 12.5, 0.4,
      size=11, italic=True, cor_txt=BRANCO, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 8 — CONCLUSÕES
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, BRANCO)
retangulo(sl, 0, 0, 13.33, 1.15, AZUL_ESC)
caixa(sl, "Conclusões", 0.4, 0.15, 12.5, 0.85,
      size=28, bold=True, cor_txt=BRANCO)
linha_h(sl, 1.25)

bullet(sl, [
    "MLP (16–8–1, tanh) atingiu R² = 0,9998 e RMSE = 0,54 kN na base sintética",
    "A rede capturou o termo não linear fc · d — algo fora do alcance de modelos lineares",
    "Normalização min-max e otimizador Adam foram essenciais para convergência estável",
    "Erros percentuais individuais inferiores a 1% em todas as 10 amostras",
    "Resultado superior ao melhor modelo da literatura recente (R² = 0,982)",
], 0.5, 1.45, 12.5, 3.5, size=19, cor_txt=CINZA, cor_bullet=AZUL_MED)

retangulo(sl, 0.4, 5.1, 12.5, 1.5, AZUL_CLA)
caixa(sl, "Perspectivas futuras", 0.6, 5.15, 4.0, 0.4,
      size=15, bold=True, cor_txt=AZUL_ESC)
bullet(sl, [
    "Ampliar base com dados experimentais reais",
    "Avaliar generalização via validação cruzada k-fold",
], 0.6, 5.55, 12.0, 0.9, size=16, cor_txt=AZUL_ESC, cor_bullet=LARANJA)

retangulo(sl, 0, 6.9, 13.33, 0.6, AZUL_MED)
caixa(sl, "III MCSM 2026  |  UNIMONTES", 0.4, 6.93, 12.5, 0.4,
      size=11, italic=True, cor_txt=BRANCO, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 9 — DISPONIBILIDADE E AGRADECIMENTOS
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, AZUL_ESC)
retangulo(sl, 0, 0, 0.18, 7.5, AZUL_MED)

caixa(sl, "Disponibilidade do Código", 0.5, 0.8, 12.5, 0.7,
      size=26, bold=True, cor_txt=BRANCO)
retangulo(sl, 0.4, 1.6, 12.5, 1.1, AZUL_MED)
caixa(sl, "github.com/JGF1987/mlp-viga-concreto-mcsm",
      0.5, 1.65, 12.3, 0.9,
      size=22, bold=True, cor_txt=BRANCO, align=PP_ALIGN.CENTER)
caixa(sl,
    "Códigos: problema4_mlp.py  |  visualizacoes_p4.py  |  gerar_word_artigo.py\n"
    "Figuras em alta resolução: pasta figuras/   |   Artigo LaTeX completo",
    0.5, 2.85, 12.5, 0.9, size=15, cor_txt=AZUL_CLA, align=PP_ALIGN.CENTER)

linha_h(sl, 4.0, AZUL_MED)

caixa(sl, "Agradecimentos", 0.5, 4.2, 12.5, 0.55,
      size=22, bold=True, cor_txt=BRANCO)
caixa(sl, "Professor Marcos e Universidade Estadual de Montes Claros — UNIMONTES",
      0.5, 4.8, 12.5, 0.6, size=18, italic=True, cor_txt=AZUL_CLA)

linha_h(sl, 5.7, AZUL_MED)

caixa(sl,
    "Maria Luiza Almeida Xavier  ·  maria.luiza.a.xavier@gmail.com\n"
    "Jordana Gabriela Fernandes   ·  jordana.funai@gmail.com\n"
    "Rayre Janaína Vieira Silva    ·  rayrejanaina2@gmail.com",
    0.5, 5.9, 12.5, 1.3, size=15, cor_txt=BRANCO)

# ── Salvar ────────────────────────────────────────────────────
out = os.path.join(PASTA, "apresentacao_mcsm2026.pptx")
prs.save(out)
print(f"Apresentação salva em: {out}")
print("Slides: Capa | Motivação | Fundamentos MLP | Dados | "
      "Métricas | Curva+Dispersão | Sensibilidade+3D | Conclusões | GitHub+Agradecimentos")
